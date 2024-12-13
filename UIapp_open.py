import streamlit as st
from langchain.chat_models import ChatOpenAI
from langchain.document_loaders import DirectoryLoader, PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.llms import Ollama
from langchain.vectorstores import FAISS
from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables.history import RunnableWithMessageHistory
from langchain_community.chat_message_histories import ChatMessageHistory
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from sentence_transformers import SentenceTransformer, util
from langchain_core.chat_history import BaseChatMessageHistory
from langchain.chains import create_history_aware_retriever
from langchain.embeddings import HuggingFaceEmbeddings
from langchain_openai import AzureOpenAIEmbeddings

import config
import langsmith
import os 
from langchain_openai import AzureChatOpenAI
from langchain.chat_models import ChatOllama
from pptx import Presentation
from langchain.chains import ConversationalRetrievalChain
# from langchain.chat_models import ChatOpenAI
from langchain.vectorstores import faiss
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.chains.conversation.memory import ConversationBufferWindowMemory
from dotenv import load_dotenv
import os


# st.set_page_config(page_title="Sales Companion Bot", layout="wide")

# Load environment variables from .env file
load_dotenv()

# Initialize LangSmith with your API key
# langsmith.api_key = "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"

# os.environ["OPENAI_API_KEY"] ="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"

os.environ["AZURE_OPENAI_API_KEY"] =os.getenv("AZURE_OPENAI_API_KEY")
os.environ["AZURE_OPENAI_ENDPOINT"] = os.getenv("AZURE_OPENAI_ENDPOINT")
os.environ["AZURE_OPENAI_API_VERSION"] = "2024-02-01"
os.environ["AZURE_OPENAI_CHAT_DEPLOYMENT_NAME"] = "gpt-35-turbo"


bot_template = '''
<div style="display: flex; align-items: center; margin-bottom: 10px;">
    <div style="flex-shrink: 0; margin-right: 10px;">
        <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
             style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
    </div>
    <div style="background-color: #dff9d8; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word; border: 1px solid #78e08f;">
        {msg}
    </div>
</div>
'''
# bot_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px;">
#     <div style="flex-shrink: 0; margin-right: 10px;">
#         <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>
#     <div style="background-color: #90EE90; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

user_template = '''
<div style="display: flex; align-items: center; margin-bottom: 10px; justify-content: flex-end;">
    <div style="flex-shrink: 0; margin-left: 10px;">
        <img src="https://cdn.iconscout.com/icon/free/png-512/free-q-characters-character-alphabet-letter-36051.png?f=webp&w=512" 
             style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
    </div>    
<div style="background-color: #ADD8E6; color: white; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
           {msg}
    </div>
</div>
'''
# def run_app():
#     llm = initialize_llm(OPENAI_API_KEY)

# Function to prepare and split documents
# @st.cache_resource
@st.cache_data
def prepare_and_split_docs(pdf_directory):
    loader = DirectoryLoader(pdf_directory, glob="**/*.pdf", loader_cls=PyPDFLoader)
    documents = loader.load()
    splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        chunk_size=1000,
        chunk_overlap=0,
        disallowed_special=(),
        separators=["\n\n", "\n", " "]
    )
    split_docs = splitter.split_documents(documents)
    return split_docs

# Function to ingest documents into vector store



def ingest_into_vectordb(split_docs):
    embedding_model1 = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
    Newdb = FAISS.from_documents(split_docs, embedding_model1)
    return Newdb



def get_conversation_chain(retriever):
    # Initialize OpenAI's ChatGPT model
    llm = AzureChatOpenAI(
        azure_deployment=os.environ["AZURE_OPENAI_CHAT_DEPLOYMENT_NAME"],  
        openai_api_version=os.environ["AZURE_OPENAI_API_VERSION"],
        temperature=0.7,
        streaming=True
    )
    
    # Contextualized question system prompt
    contextualize_q_system_prompt = (
        "You are a helpful assistant. Using the chat history and the latest user question, "
        "provide a relevant response based on the retrieved documents. If no relevant answer is found, "
        "respond with: 'I'm sorry, but I couldn’t find an answer. Could you rephrase or provide more details?'"
    )
    
    contextualize_q_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", contextualize_q_system_prompt),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}"),
        ]
    )

    history_aware_retriever = create_history_aware_retriever(
        llm, retriever, contextualize_q_prompt
    )

    # System prompt for the question-answering chain
    system_prompt = (
        "As a chat assistant, your task is to provide accurate and concise information based on the provided documents. "
        "For each query, identify relevant information and respond in 2-3 sentences. If no relevant information is found, "
        "reply with: 'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?' "
        "{context}"
    )

    qa_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", system_prompt),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}"),
        ]
    )

    question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)
    rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)

    store = {}

    def get_session_history(session_id: str) -> BaseChatMessageHistory:
        if session_id not in store:
            store[session_id] = ChatMessageHistory()
        return store[session_id]

    conversational_rag_chain = RunnableWithMessageHistory(
        rag_chain,
        get_session_history,
        input_messages_key="input",
        history_messages_key="chat_history",
        output_messages_key="answer",
    )
    
    return conversational_rag_chain

def get_relevant_context(query, retriever):
    docs = retriever.get_relevant_documents(query)
    return [doc[:500] for doc in docs]  # Truncate to 500 characters


pdf_directory="Data/SalesDocs"
split_docs = prepare_and_split_docs(pdf_directory)
vector_db = ingest_into_vectordb(split_docs)
retriever = vector_db.as_retriever()

#  Function to create PPT from chat history
def create_ppt(chat_history, title="Chatbot Responses", subtitle="Generated Presentation"):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    for i, message in enumerate(chat_history):
        slide_layout = prs.slide_layouts[1]  # Title and content slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = f"Response {i + 1}"
        content.text = message['bot']

    prs.save('chatbot_responses.pptx')


# Main Streamlit app
# st.markdown(css, unsafe_allow_html=True)
col1, col2 = st.columns([1, 1])
logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png"
import streamlit as st

# Set the page configuration
# st.set_page_config(page_title="Sales Companion Bot", layout="wide")

# Apply custom CSS for better styling
st.markdown("""
<style>
    body {
        background-color: #f0f4f8;  /* Light background color */
    }
    .chat-container {
        border-radius: 10px;
        padding: 20px;
        background-color: #ffffff;  /* White background for chat */
        box-shadow: 0 4px 10px rgba(0, 0, 0, 0.1);
        max-width: 600px;
        margin: auto;
    }
    .user-message {
        background-color: #90EE90;  /* Light green for user messages */
        padding: 10px;
        border-radius: 10px;
        margin: 5px 0;
        text-align: left;
        border: 1px solid #78e08f;
    }
    .bot-message {
        background-color: #dff9d8;  /* Light orange for bot messages */
        padding: 10px;
        border-radius: 10px;
        margin: 5px 0;
        text-align: left;
        border: 1px solid #78e08f;
    }
</style>
""", unsafe_allow_html=True)


st.title("Sales Companion Bot")
st.write("Welcome! How can I help you today?")

pdf_directory = "Data/SalesDocs"
split_docs = prepare_and_split_docs(pdf_directory)
vector_db = ingest_into_vectordb(split_docs)
retriever = vector_db.as_retriever()
conversational_chain = get_conversation_chain(retriever)
st.session_state.conversational_chain = conversational_chain

# Initialize chat history if not present
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []

# Chat input area
user_input = st.text_input("Ask a question about the provided documents:")

# Create columns for layout
col1, col2 = st.columns([1, 1])

with col1:
    if st.button("Submit"):
        user_input_lower = user_input.lower()  # Normalize input for checking
        st.session_state.user_input = user_input

        # Respond to greetings
        if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
            response = "Hello! How can I assist you?"
            st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

        # Respond to thanks
        elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
            response = "Thank you! Let me know if you have any more queries."
            st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

        # Respond to yes/no
        elif any(yn in user_input_lower for yn in ["yes", "no"]):
            response = "I understand. If you have any other questions, feel free to ask!"
            st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

        # Process document queries
        elif user_input and 'conversational_chain' in st.session_state:
            session_id = "user412"  # Static session ID for this demo; you can make it dynamic if needed
            
            # # Check for specific queries
            # if "meeting a new prospect" in user_input_lower or "latest corporate presentation" in user_input_lower:
            #     response = {
            #         "answer": "Please find below the document for your query.",
            #         "context_docs": split_docs  # Show all documents related to the query
            #     }
            # else:
            #     response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
             # Check for specific queries
if "meeting a new prospect" in user_input_lower or "latest corporate presentation" in user_input_lower:
    # Create a response with document names and paths
    response = {
        "answer": "Please find below the document for your query.",
        "context_docs": [
            {"name": doc, "path": f"{pdf_directory}/{doc}"} for doc in split_docs  # Assuming split_docs contains the document names
        ]
    }
else:
    response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})   
            # Insert the new response at the top of the chat history
st.session_state.chat_history.insert(0, {"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})

with col2:
    if st.button("Refresh"):
        st.session_state.user_input = ""  # Clear the input
        st.session_state.chat_history = []  # Clear chat history

# Display chat history
if st.session_state.chat_history:
    for message in st.session_state.chat_history:
        st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
        st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

# Sidebar for frequently asked questions
with st.sidebar:
    st.subheader("Frequently Asked Questions")
    session_id = "user123"
    for message in st.session_state.chat_history:
        if st.button(message['user']):
            st.session_state.selected_question = message['user']
            st.session_state.selected_answer = message['bot']
# ----------------------------------------------------------
# st.image(logopath, width=200)
# st.title("Sales Companion Bot")
# st.write("Welcome! How can I help you today?")

# pdf_directory = "Data/SalesDocs"  # Update to your actual folder path

# # Prepare documents and ingest into vector store
# split_docs = prepare_and_split_docs(pdf_directory)
# vector_db = ingest_into_vectordb(split_docs)
# retriever = vector_db.as_retriever()

# # Initialize conversation chain
# conversational_chain = get_conversation_chain(retriever)
# st.session_state.conversational_chain = conversational_chain

# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []

# # Chat input
# # conversation_chain = initialize_conversation_chain(llm)
# user_input = st.text_input("Ask a question about the provided documents:")


# col1, col2 = st.columns([1, 1])

# with col1:
#     if st.button("Submit"):
#         user_input_lower = user_input.lower()  # Normalize input for checking
        

#         # Store the user input in session state
#         st.session_state.user_input = user_input
#         # Respond to greetings
#         if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
#             response = "Hello! How can I assist you?"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})
           

#         # Respond to thanks
#         elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
#             response = "Thank you! Let me know if you have any more queries."
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to yes/no
#         elif any(yn in user_input_lower for yn in ["yes", "no"]):
#             response = "I understand. If you have any other questions, feel free to ask!"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Process document queries
#         elif user_input and 'conversational_chain' in st.session_state:
#             session_id ="user412" # Static session ID for this demo; you can make it dynamic if needed
#             response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#             # context_docs = response.get('context', [])


        
#             # Check for specific queries
#         if "meeting a new prospect" in user_input_lower or "latest corporate presentation" in user_input_lower:
#                 response = {
#                     "answer": "Please find below the document for your query.",
#                     "context_docs": split_docs  # Assuming you want to show all documents related to the query
#                 }
#         else:
#                 session_id ="user412"
#                 response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})

#                 context_docs = response.get('context', [])
                
#                  # Insert the new response at the top of the chat history
        
#         st.session_state.chat_history.insert(0, {"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})
#         # Append the chat history
#         st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})

#     # with col2:
#     #  if st.button("Refresh"):
#     #     st.session_state.user_input = ""  # Clear the input

# if user_input and 'conversational_chain' in st.session_state:
#    session_id ="user412" # Static session ID for this demo; you can make it dynamic if needed
#    response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#    context_docs = response.get('context', [])

# # st.session_state.chat_history.insert(0, {"user": user_input, "bot": response})

# with st.sidebar:
#     st.subheader("Frequently Asked Questions")
#     session_id ="user123"
#     for message in st.session_state.chat_history:
#         if st.button(message['user']):
#             st.session_state.selected_question = message['user']
#             st.session_state.selected_answer = message['bot']

# # Display chat history
# if 'selected_question' in st.session_state and 'selected_answer' in st.session_state:
#     st.markdown(f"<div style='text-align: left;'><strong>You:</strong> {st.session_state.selected_question}</div>", unsafe_allow_html=True)
#     st.markdown(f"<div style='background-color: #dff9d8; padding: 10px; border-radius: 10px; border: 1px solid #78e08f;'><strong>Bot:</strong> {st.session_state.selected_answer}</div>", unsafe_allow_html=True)

# Clear history and download options
# col2 = st.columns(2)

# with col2[1]:
#     if st.button("Refresh"):
#         st.session_state.user_input = ""  # Clear the input
#         st.session_state.chat_history = []  # Clear chat 
#     if st.button("Download PPT"):
#         if st.session_state.chat_history:
#             create_ppt(st.session_state.chat_history)
#             st.success("PPT created! You can download it [here](chatbot_responses.pptx).")
#         else:
#             st.warning("No chat history available to create a PPT.")


   
# with col2:
#     if st.button("Refresh"):
#         st.session_state.user_input = ""  # Clear the input
#         st.session_state.chat_history = []  # Clear chat 
#     if st.button("Download PPT"):
#         if st.session_state.chat_history:
#             create_ppt(st.session_state.chat_history)
#             st.success("PPT created! You can download it [here](chatbot_responses.pptx).")
#         else:
#             st.warning("No chat history available to create a PPT.")

#          # Display chat history
# # Display chat history
# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

#         if message.get("context_docs"):
#             with st.expander("Source Documents"):
#                 for doc in message["context_docs"]:
#                     st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
#                     st.write(doc.page_content)



# import streamlit as st
# from langchain.chat_models import ChatOpenAI, AzureChatOpenAI
# from langchain.document_loaders import DirectoryLoader, PyPDFLoader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain.vectorstores import FAISS
# from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
# from langchain_core.runnables.history import RunnableWithMessageHistory
# from langchain_community.chat_message_histories import ChatMessageHistory
# from langchain.chains import create_retrieval_chain, create_stuff_documents_chain
# from langchain.embeddings import HuggingFaceEmbeddings
# from dotenv import load_dotenv
# import os
# from pptx import Presentation

# Load environment variables from .env file
# load_dotenv()

# # Set up Azure OpenAI API keys
# os.environ["AZURE_OPENAI_API_KEY"] = os.getenv("AZURE_OPENAI_API_KEY")
# os.environ["AZURE_OPENAI_ENDPOINT"] = os.getenv("AZURE_OPENAI_ENDPOINT")
# os.environ["AZURE_OPENAI_API_VERSION"] = "2024-02-01"
# os.environ["AZURE_OPENAI_CHAT_DEPLOYMENT_NAME"] = "gpt-35-turbo"

# # Templates for bot and user messages
# bot_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px;">
#     <div style="flex-shrink: 0; margin-right: 10px;">
#         <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>
#     <div style="background-color: #FFA07A; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word; border: 1px solid #f08080;">
#         {msg}
#     </div>
# </div>
# '''

# user_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px; justify-content: flex-end;">
#     <div style="flex-shrink: 0; margin-left: 10px;">
#         <img src="https://cdn.iconscout.com/icon/free/png-512/free-q-characters-character-alphabet-letter-36051.png?f=webp&w=512" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>    
#     <div style="background-color: #808080; color: white; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#            {msg}
#     </div>
# </div>
# '''
# @st.cache_data
# def prepare_and_split_docs(pdf_directory):
#     loader = DirectoryLoader(pdf_directory, glob="**/*.pdf", loader_cls=PyPDFLoader)
#     documents = loader.load()
#     splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
#         chunk_size=1000,
#         chunk_overlap=0
#     )
#     split_docs = splitter.split_documents(documents)
#     return split_docs

# def ingest_into_vectordb(split_docs):
#     embedding_model = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
#     vector_db = FAISS.from_documents(split_docs, embedding_model)
#     return vector_db

# def get_conversation_chain(retriever):
#     llm = AzureChatOpenAI(
#         azure_deployment=os.environ["AZURE_OPENAI_CHAT_DEPLOYMENT_NAME"],
#         openai_api_version=os.environ["AZURE_OPENAI_API_VERSION"],
#         temperature=0.7,
#         streaming=True
#     )

#     contextualize_q_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", "You are a helpful assistant."),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )

#     history_aware_retriever = create_history_aware_retriever(
#         llm, retriever, contextualize_q_prompt
#     )

#     system_prompt = "As a chat assistant, provide accurate responses based on the documents. If no answer is found, respond with: 'I'm sorry, but I couldn’t find an answer.' {context}"

#     qa_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )

#     question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)
#     rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)

#     store = {}
#     def get_session_history(session_id: str):
#         if session_id not in store:
#             store[session_id] = ChatMessageHistory()
#         return store[session_id]

#     conversational_rag_chain = RunnableWithMessageHistory(
#         rag_chain,
#         get_session_history,
#         input_messages_key="input",
#         history_messages_key="chat_history",
#         output_messages_key="answer",
#     )
    
#     return conversational_rag_chain
# # Initialize chat history if it doesn't exist
# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []



# def create_ppt(chat_history, title="Chatbot Responses"):
#     prs = Presentation()
#     slide_layout = prs.slide_layouts[0]  # Title slide
#     slide = prs.slides.add_slide(slide_layout)
#     slide.shapes.title.text = title

#     for i, message in enumerate(chat_history):
#         slide_layout = prs.slide_layouts[1]  # Title and content slide
#         slide = prs.slides.add_slide(slide_layout)
#         slide.shapes.title.text = f"Response {i + 1}"
#         slide.placeholders[1].text = message['bot']

#     prs.save('chatbot_responses.pptx')

# # Main Streamlit app
# st.title("Sales Companion Bot")
# st.write("Welcome! How can I help you today?")

# pdf_directory = "Data/SalesDocs"
# split_docs = prepare_and_split_docs(pdf_directory)
# vector_db = ingest_into_vectordb(split_docs)
# retriever = vector_db.as_retriever()
# conversational_chain = get_conversation_chain(retriever)
# st.session_state.conversational_chain = conversational_chain

# # Initialize chat history if it doesn't exist
# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []


# # Initialize chat history if not present
# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []

# # Chat input area
# user_input = st.text_input("Ask a question about the provided documents:")

# # Create columns for layout
# col1, col2 = st.columns([1, 1])

# with col1:
#     if st.button("Submit"):
#         user_input_lower = user_input.lower()  # Normalize input for checking
#         st.session_state.user_input = user_input

#         # Respond to greetings
#         if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
#             response = "Hello! How can I assist you?"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to thanks
#         elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
#             response = "Thank you! Let me know if you have any more queries."
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to yes/no
#         elif any(yn in user_input_lower for yn in ["yes", "no"]):
#             response = "I understand. If you have any other questions, feel free to ask!"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Process document queries
#         elif user_input and 'conversational_chain' in st.session_state:
#             session_id = "user412"  # Static session ID for this demo; you can make it dynamic if needed
            
#             # Check for specific queries
#             if "meeting a new prospect" in user_input_lower or "latest corporate presentation" in user_input_lower:
#                 response = {
#                     "answer": "Please find below the document for your query.",
#                     "context_docs": split_docs  # Show all documents related to the query
#                 }
#             else:
#                 response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
                
#             # Insert the new response at the top of the chat history
#             st.session_state.chat_history.insert(0, {"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})
# # Initialize chat history if it doesn't exist
# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []


# with col2:
#     if st.button("Refresh"):
#         st.session_state.user_input = ""  # Clear the input
#         st.session_state.chat_history = []  # Clear chat history

# # Display chat history
# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

# # Sidebar for frequently asked questions
# with st.sidebar:
#     st.subheader("Frequently Asked Questions")
#     session_id = "user123"
#     for message in st.session_state.chat_history:
#         if st.button(message['user']):
#             st.session_state.selected_question = message['user']
#             st.session_state.selected_answer = message['bot']