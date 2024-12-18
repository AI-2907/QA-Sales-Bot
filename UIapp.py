
# =================================code without Next quetion button========================
import streamlit as st
from langchain.chat_models import ChatOpenAI
from langchain.document_loaders import DirectoryLoader, PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.llms import Ollama
from langchain.vectorstores import FAISS
from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables.history import RunnableWithMessageHistory
from langchain_community.chat_message_histories import ChatMessageHistory
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from sentence_transformers import SentenceTransformer, util
from langchain_core.chat_history import BaseChatMessageHistory
from langchain.chains import create_history_aware_retriever
from langchain.embeddings import HuggingFaceEmbeddings
from langchain_openai import AzureOpenAIEmbeddings

import config
import langsmith
import os 
from langchain_openai import AzureChatOpenAI
from langchain.chat_models import ChatOllama
from pptx import Presentation
from langchain.chains import ConversationalRetrievalChain
# from langchain.chat_models import ChatOpenAI
from langchain.vectorstores import faiss
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.chains.conversation.memory import ConversationBufferWindowMemory
from dotenv import load_dotenv
import os

# Load environment variables from .env file
load_dotenv()

# os.environ["OPENAI_API_KEY"] ="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"

os.environ["AZURE_OPENAI_API_KEY"] =os.getenv("AZURE_OPENAI_API_KEY")
os.environ["AZURE_OPENAI_ENDPOINT"] = os.getenv("AZURE_OPENAI_ENDPOINT")
os.environ["AZURE_OPENAI_API_VERSION"] = "2024-02-01"
os.environ["AZURE_OPENAI_CHAT_DEPLOYMENT_NAME"] = "gpt-35-turbo"

# Templates for user and bot messages
bot_template = '''
<div style="display: flex; align-items: center; margin-bottom: 10px;">
    <div style="flex-shrink: 0; margin-right: 10px;">
        <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
             style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
    </div>
    <div style="background-color: #dff9d8;color: black; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word; border: 1px solid #78e08f;">
        {msg}
    </div>
</div>
'''
# bot_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px;">
#     <div style="flex-shrink: 0; margin-right: 10px;">
#         <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>
#     <div style="background-color: #f1f1f1; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

user_template = '''
<div style="display: flex; align-items: center; margin-bottom: 10px; justify-content: flex-end;">
    <div style="flex-shrink: 0; margin-left: 10px;">
        <img src="https://cdn.iconscout.com/icon/free/png-512/free-q-characters-character-alphabet-letter-36051.png?f=webp&w=512" 
             style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
    </div>    
    <div style="background-color: #66bfff; color: black; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
        {msg}
    </div>
</div>
'''
# def run_app():
#     llm = initialize_llm(OPENAI_API_KEY)

# Function to prepare and split documents
# @st.cache_resource
@st.cache_data
def prepare_and_split_docs(pdf_directory):
    loader = DirectoryLoader(pdf_directory, glob="**/*.pdf", loader_cls=PyPDFLoader)
    documents = loader.load()
    splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        chunk_size=512,
        chunk_overlap=256,
        disallowed_special=(),
        separators=["\n\n", "\n", " "]
    )
    split_docs = splitter.split_documents(documents)
    return split_docs

# Function to ingest documents into vector store

# def ingest_into_vectordb(split_docs):
#     # Use OpenAI embeddings with `text-embedding-ada-002`
#     embeddings=AzureOpenAIEmbeddings(model="text-embedding-3-large")
#     # embeddings = OpenAIEmbeddings(model="text-embedding-ada-002", api_key="AZURE_OPENAI_API_KEY")
    
#     # Create a FAISS vectorstore from the documents
#     db = FAISS.from_documents(split_docs, embeddings)
    
#     # Save the vectorstore locally
#     DB_FAISS_PATH = 'vectorstore/db_faiss'
#     db.save_local(DB_FAISS_PATH)
    
#     print("Documents are inserted into FAISS vectorstore")
#     return db

def ingest_into_vectordb(split_docs):
    embedding_model1 = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
    Newdb = FAISS.from_documents(split_docs, embedding_model1)
    return Newdb


def get_conversation_chain(retriever):
    # Initialize OpenAI's ChatGPT model
    llm = AzureChatOpenAI(
        azure_deployment=os.environ["AZURE_OPENAI_CHAT_DEPLOYMENT_NAME"],  
        openai_api_version=os.environ["AZURE_OPENAI_API_VERSION"],
        temperature=0.9,
        streaming=True 
        
     )
    contextualize_q_system_prompt = (
        "Given the chat history and the latest user question, "
        "provide a response based on the documents. If no answer is found, "
        "respond: 'I'm sorry, but I couldn’t find an answer. Could you rephrase or provide more details?'"
    )
    contextualize_q_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", contextualize_q_system_prompt),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}"),
        ]
    )
    history_aware_retriever = create_history_aware_retriever(
        llm, retriever, contextualize_q_prompt
    )
    system_prompt = (
        "As a chat assistant, provide accurate and relevant information  based on the provided document in 2-3 sentences.accoring to given query find matching wordsfrom query in provided documents and give response using that "
        "Answer should be correct to the point short and brief for given quetion . If no relevant information is found, respond with: "
        "'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?' "
        "{context}"
    )
    qa_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", system_prompt),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}"),
        ]
    )
    question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)
    rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)

    history_aware_retriever = create_history_aware_retriever(
        llm, retriever, contextualize_q_prompt
    )
    qa_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", "Answer concisely based on the documents. {context}"),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}"),
        ]
    )
    question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)
    rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)

    store = {}

    def get_session_history(session_id: str) -> BaseChatMessageHistory:
        if session_id not in store:
            store[session_id] = ChatMessageHistory()
        return store[session_id]

    conversational_rag_chain = RunnableWithMessageHistory(
        rag_chain,
        get_session_history,
        input_messages_key="input",
        history_messages_key="chat_history",
        output_messages_key="answer",
    )
    return conversational_rag_chain

def get_relevant_context(query, retriever):
    docs = retriever.get_relevant_documents(query)
    return [doc[:500] for doc in docs]  # Truncate to 500 characters


pdf_directory="Data/SalesDocs"
split_docs = prepare_and_split_docs(pdf_directory)
vector_db = ingest_into_vectordb(split_docs)
retriever = vector_db.as_retriever()

# def initialize_conversation_chain(llm):
#     return ConversationalRetrievalChain(
#         llm=llm,
#         memory=ConversationBufferWindowMemory(k=5)  # Stores the last 5 interactions.
#     )

#  Function to create PPT from chat history
# def create_ppt(chat_history, title="Chatbot Responses", subtitle="Generated Presentation"):
#     prs = Presentation()
#     slide_layout = prs.slide_layouts[0]  # Title slide
#     slide = prs.slides.add_slide(slide_layout)
#     title_placeholder = slide.shapes.title
#     subtitle_placeholder = slide.placeholders[1]
#     title_placeholder.text = title
#     subtitle_placeholder.text = subtitle

#     for i, message in enumerate(chat_history):
#         slide_layout = prs.slide_layouts[1]  # Title and content slide
#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         content = slide.placeholders[1]
        
#         title.text = f"Response {i + 1}"
#         content.text = message['bot']

#     prs.save('chatbot_responses.pptx')


# Main Streamlit app
# col1, col2 = st.columns([1, 1])
# logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png

# Set the page configuration
# st.set_page_config(page_title="Sales Companion Bot", layout="wide")
# Load the CSS file
with open("style.css") as f:
    st.markdown(f'<style>{f.read()}</style>', unsafe_allow_html=True)

# Apply custom CSS for better styling
st.markdown("""
<style>
    body {
        background-color: #f0f4f8;  /* Light background color for the whole app */
    }
    .chat-container {
        border-radius: 10px;
        padding: 20px;
        background-color: #ADD8E6;  /* White background for chat */
        box-shadow: 0 4px 10px rgba(0, 0, 0, 0.1);
        max-width: 600px;
        margin: auto;
    }
    .user-message {
        background-color: #90EE90;  /* Light green for user messages */
        padding: 10px;
        border-radius: 10px;
        margin: 5px 0;
        text-align: left;
        border: 1px solid #78e08f;
    }
    .bot-message {
        background-color: #dff9d8;  /* Light orange for bot messages */
        padding: 10px;
        border-radius: 10px;
        margin: 5px 0;
        text-align: left;
        border: 1px solid #78e08f;
    }
    .sidebar .sidebar-content {
        background-color: #e0f7fa;  /* Light cyan for sidebar */
        padding: 10px;
        border-radius: 10px;
    }
</style>
""", unsafe_allow_html=True)

# Initialize chat history if it doesn't exist
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
 # Initialize conversation chain
conversational_chain = get_conversation_chain(retriever)
st.session_state.conversational_chain = conversational_chain
# Set up the app title and welcome message
logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png"
st.image(logopath, width=200)
st.title("Sales Companion Bot")
st.write("Welcome! How can I help you today?")
# Recommended prompts
recommended_prompts = [
    "Summarize case studies for Data Analytics",
    "Tell me the top 3 case studies",
]

# Sidebar for recommended prompts
with st.sidebar:
    st.markdown("<div class='box'><h4>Recommended Prompts</h4></div>", unsafe_allow_html=True)
    for prompt in recommended_prompts:
        if st.button(prompt):
            st.session_state.user_input = prompt
            # Simulating a response for the clicked prompt
            if prompt == "Summarize case studies for Data Analytics":
                # response = {"answer": "Here are some key insights from the case studies in Data Analytics...", "context_docs": []}
                session_id ="user412"
                response = st.session_state.conversational_chain.invoke({"input": prompt}, config={"configurable": {"session_id": session_id}})
                context_docs = response.get('context', [])
            elif prompt == "Tell me the top 3 case studies":
                # response = {"answer": "1. Case Study A\n2. Case Study B\n3. Case Study C", "context_docs": []}
                session_id ="user412"
                response = st.session_state.conversational_chain.invoke({"input": prompt}, config={"configurable": {"session_id": session_id}})
                context_docs = response.get('context', [])
            
            # Add the response to chat history
            st.session_state.chat_history.append({
                "user": prompt,
                "bot": response.get('answer', "I didn't understand that."),
                "context_docs": response.get("context_docs", [])
            })
# Buttons for Hunting and Farming
col1, col2 = st.columns(2)

with col1:
    if st.button("Hunting"):
        st.session_state.mode = "hunting"
        st.write("You selected Hunting mode.")

with col2:
    if st.button("Farming"):
        st.session_state.mode = "farming"
        st.write("You selected Farming mode.")

# Display a textbox for prompts when Farming mode is selected
if st.session_state.get('mode') == "farming":
    user_input = st.text_input("Enter your prompt for Farming:")

    if st.button("Submit"):
        user_input_lower = user_input.lower()  # Normalize input for checking

        # Store the user input in session state
        st.session_state.user_input = user_input

        # Respond to greetings
        if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
            response = "Hello! How can I assist you?"
        
        # Respond to thanks
        elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
            response = "Thank you! Let me know if you have any more queries."
        
        # Respond to yes/no
        elif any(yn in user_input_lower for yn in ["yes", "no"]):
            response = "I understand. If you have any other questions, feel free to ask!"
        
        # Process document queries
        elif user_input and 'conversational_chain' in st.session_state:
            session_id = "user412"  # Static session ID for this demo; you can make it dynamic if needed
            response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})

              # Check for specific queries
        if "I am meeting a new prospect" in user_input_lower or "latest corporate presentation" in user_input_lower:
                response = {
                    "answer": "Please find below the document for your query.",
                    "context_docs": split_docs  # Assuming you want to show all documents related to the query
                }
            
        else:
                session_id ="user412"
                response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
                context_docs = response.get('context', [])
        #          # Insert the new response at the top of the chat history
        st.session_state.chat_history.insert(0, {"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})
        response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
        # if isinstance(response, str):
        #             # If response is a string, wrap it in a dictionary
        #             response = {"answer": response, "context_docs": []}

        # # Storing the response in chat history
        # st.session_state.chat_history.append({
        #     "user": user_input,
        #     # "bot": response.get('answer', "I didn't understand that."),
        #     "context_docs": response.get("context_docs", [])
        # })

if st.session_state.chat_history:
    for message in st.session_state.chat_history:
        st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
        st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

        if message.get("context_docs"):
            with st.expander("Source Documents"):
                for doc in message["context_docs"]:
                    st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
                    st.write(doc.page_content)       
# # Chat history display
# if st.session_state.chat_history:
#     chat_container = st.container()
#     with chat_container:
#         for message in st.session_state.chat_history:
#             st.markdown(f"<div class='user-message'>{message['user']}</div>", unsafe_allow_html=True)
#             st.markdown(f"<div class='bot-message'>{message['bot']}</div>", unsafe_allow_html=True)

# # Chat history display
# if st.session_state.chat_history:
#     chat_container = st.container()
#     with chat_container:
#         for message in st.session_state.chat_history:
#             st.markdown(f"<div class='user-message'>{message['user']}</div>", unsafe_allow_html=True)
            # st.markdown(f"<div class='bot-message'>{message['bot']}</div>", unsafe_allow_html=True)

# # Sidebar for frequently asked questions
# with st.sidebar:
#     st.subheader("Frequently Asked Questions")
#     st.markdown("<div class='sidebar-content'>", unsafe_allow_html=True)  # Sidebar content styling
#     for message in st.session_state.chat_history:
#         if st.button(message['user']):
#             st.session_state.selected_question = message['user']
#             st.session_state.selected_answer = message['bot']
#     st.markdown("</div>", unsafe_allow_html=True)  # Close sidebar content styling
# st.image(logopath, width=200)
# st.title("Sales Companion Bot")
# st.write("Welcome! How can I help you today?")

# pdf_directory = "Data/SalesDocs"  # Update to your actual folder path

# # Prepare documents and ingest into vector store
# split_docs = prepare_and_split_docs(pdf_directory)
# vector_db = ingest_into_vectordb(split_docs)
# retriever = vector_db.as_retriever()

# # Initialize conversation chain
# conversational_chain = get_conversation_chain(retriever)
# st.session_state.conversational_chain = conversational_chain

# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []

# # Chat input
# # conversation_chain = initialize_conversation_chain(llm)
# user_input = st.text_input("Ask a question about the provided documents:")
#

# Chat history display
# if st.session_state.chat_history:
#     chat_container = st.container()
#     with chat_container:
#         for message in st.session_state.chat_history:
#             st.markdown(f"<div class='user-message'>{message['user']}</div>", unsafe_allow_html=True)
#             st.markdown(f"<div class='bot-message'>{message['bot']}</div>", unsafe_allow_html=True)

# Sidebar for frequently asked questions
# with st.sidebar:
    # st.subheader("Frequently Asked Questions")
    # st.markdown("<div class='sidebar-content'>", unsafe_allow_html=True)  # Sidebar content styling
    # for message in st.session_state.chat_history:
    #     if st.button(message['user']):
    #         st.session_state.selected_question = message['user']
    #         st.session_state.selected_answer = message['bot']
    # st.markdown("</div>", unsafe_allow_html=True)  # Close sidebar content styling

# col1, col2 = st.columns([1, 1])
# +++++++++++++++++++
# with col1:
#     if st.button("Submit"):
#         user_input_lower = user_input.lower()  # Normalize input for checking
        

#     #     # Store the user input in session state
#     #     st.session_state.user_input = user_input
#     #     # Respond to greetings
#         if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
#             response = "Hello! How can I assist you?"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to thanks
#         elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
#             response = "Thank you! Let me know if you have any more queries."
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to yes/no
#         elif any(yn in user_input_lower for yn in ["yes", "no"]):
#             response = "I understand. If you have any other questions, feel free to ask!"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Process document queries
#         elif user_input and 'conversational_chain' in st.session_state:
#             session_id ="user412" # Static session ID for this demo; you can make it dynamic if needed
#             response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#             # context_docs = response.get('context', [])


#         # Process document queries
#         # elif user_input and 'conversational_chain' in st.session_state:
#         #     session_id = "user412"  # Static session ID for this demo; you can make it dynamic if needed
#         #     response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})

        #     # Check for specific queries
        # if "meeting a new prospect" in user_input_lower or "latest corporate presentation" in user_input_lower:
        #         response = {
        #             "answer": "Please find below the document for your query.",
        #             "context_docs": split_docs  # Assuming you want to show all documents related to the query
        #         }
#         else:
#                 session_id ="user412"
#                 response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#                 context_docs = response.get('context', [])
#         #          # Insert the new response at the top of the chat history
#         st.session_state.chat_history.insert(0, {"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})
#         # Check for specific queries
#         # Check for specific queries
#         session_id ="user412"
#         if "meeting a new prospect" in user_input_lower or "latest corporate presentation" in user_input_lower:
#             session_id ="user412"
#     # Create a response with specific document names and paths
#             response = {
#         "answer": "Please find below the document for your query.",
#         "context_docs": [
#             {"name": "Meeting_New_Prospect.pdf", "path": f"{pdf_directory}/Meeting_New_Prospect.pdf"},
#             {"name": "Latest_Corporate_Presentation.pdf", "path": f"{pdf_directory}/Latest_Corporate_Presentation.pdf"}
#         ]  # Only include relevant documents
#     }
#     else:
#         session_id ="user412"
#         response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
# ++++++++++++++++++++++++/


            # Append the chat history
        # st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})

# col1, col2 = st.columns([1, 1])
# with col1:
#     if st.button("Submit"):
#         user_input_lower = user_input.lower()  # Normalize input for checking

#         # Store the user input in session state
#         st.session_state.user_input = user_input
        

#         # Respond to greetings
#         if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
#             response = "Hello! How can I assist you?"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to thanks
#         elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
#             response = "Thank you! Let me know if you have any more queries."
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to yes/no
#         elif any(yn in user_input_lower for yn in ["yes", "no"]):
#             response = "I understand. If you have any other questions, feel free to ask!"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Process document queries
#         elif user_input and 'conversational_chain' in st.session_state:
#             session_id ="user412" # Static session ID for this demo; you can make it dynamic if needed
#             response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#             context_docs = response.get('context', [])

#             # Check if the response is the default message
#             if response['answer'] == "I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?":
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": []})
#                 # Check for specific queries
#                 if "meeting a new prospect" in user_input_lower or "latest corporate presentation" in user_input_lower:
#                  response = {
#                     "answer": "Please find below the document for your query.",
#                     "context_docs": split_docs  # Assuming you want to show all documents related to the query
#                 }
#             else:
#                 response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#                 context_docs = response.get('context', [])

#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": context_docs})
#                 st.session_state.chat_history.insert(0, {"role": "assistant", "content": context_docs})
#         # Use LangSmith for logging
#         # with langsmith.log_session("SalesBot"):
#         #     response = conversational_chain.invoke({"input": user_input}, config={"session_id": "user123"})
#         #     context_docs = response.get("context", [])
#         #     response = response["answer"]
#         #     if response == "I'm sorry, but I couldn’t find an answer.":
#         #         context_docs = []

#         # # Log response and context to LangSmith
#         # langsmith.log_event("Response", data={"response": response})
#         # langsmith.log_event("Context Docs", data={"docs": context_docs})

#     # st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": context_docs})


# with col2:
#     if st.button("Refresh"):
#         st.session_state.user_input = ""  # Clear the input
#         st.session_state.chat_history = []  # Clear chat 
#     # if st.button("Download PPT"):
#         if st.session_state.chat_history:
        #     create_ppt(st.session_state.chat_history)
        #     st.success("PPT created! You can download it [here](chatbot_responses.pptx).")
        # else:
        #     st.warning("No chat history available to create a PPT.")
# with col2:
#     if st.button("Refresh"):
#         st.session_state.user_input = ""  # Clear the input

# col1, col2 = st.columns([1, 1])
# if st.button("Submit"):
#     user_input_lower = user_input.lower()
#     greetings = ["hi", "hello", "hey"]
#     thanks = ["thank you", "thanks", "thx"]

#     if any(greet in user_input_lower for greet in greetings):
#         response = "Hello! How can I assist you?"
#     elif any(thank in user_input_lower for thank in thanks):
#         response = "Thank you! Let me know if you have more queries."
#     else:
#         response = conversational_chain.invoke({"input": user_input}, config={"session_id": "user123"})
#         context_docs = response.get("context", [])
#         response = response["answer"]
#         if response == "I'm sorry, but I couldn’t find an answer.":
#             context_docs = []

#     st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": context_docs})
# with col2:
#     if st.button("Next Question"):
#         st.session_state.user_input = ""
# response=""
# if st.session_state.chat_history:
#  for message in st.session_state.chat_history:
# # for message in reversed(st.session_state.chat_history):
#     user_msg = message.get("user")
#     bot_msg = message.get("bot")

    # st.markdown(user_template.format(msg=user_msg), unsafe_allow_html=True)
    # st.markdown(bot_template.format(msg=bot_msg), unsafe_allow_html=True)
# Display chat history
# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)
# # for message in st.session_state.chat_history:
#     st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#     st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

# ++++++++++++++++++++++++
# Display chat history/
# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

#         if message.get("context_docs"):
#             with st.expander("Source Documents"):
#                 for doc in message["context_docs"]:
#                     st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
#                     st.write(doc.page_content)
# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         if 'user' in message and 'bot' in message:
#             st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#             st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)
#         else:
#             st.error("Invalid message format: missing 'user' or 'bot' key.")


#     if message.get("context_docs"):
#         with st.expander("Source Documents"):
#             for doc in message["context_docs"]:
#                 st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
#                 st.write(doc.page_content)


# Display chat history (newest at the top)
    # st.write("### Chat History")
    # for message in st.session_state.messages:
    #     if message["role"] == "user":
    #         st.markdown(f"**You:** {message['content']}")
    #     else:
    #         st.markdown(f"**Bot:** {message['content']}")

# if __name__ == "__main__":
#     main()

def display_chat_history(session_id):
    history = store.get(session_id, ChatMessageHistory())
    for msg in history.messages:
        if msg.type == "human":
            st.write(f"**You:** {msg.content}")
        elif msg.type == "system":
            st.write(f"**Bot:** {msg.content}")
# +++++++++++++++++++++++++++++++++++++++++++++++
# import streamlit as st
# # from your_module import prepare_and_split_docs, ingest_into_vectordb, get_conversation_chain

# # Display the logo and welcome message
# logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png"
# st.image(logopath, width=200)
# st.title("Q&A Bot")
# st.write("Welcome! How can I help you today?")

# # Define the PDF directory (update path as per your environment)
# pdf_directory = "Data/SalesDocs"

# # Prepare documents and ingest them into vector storage
# try:
#     st.write("Preparing documents for search...")
#     split_docs = prepare_and_split_docs(pdf_directory)  # Function to read and split PDFs
#     vector_db = ingest_into_vectordb(split_docs)       # Function to ingest documents into a vector DB
#     retriever = vector_db.as_retriever()               # Retrieve documents using a retriever interface

#     # Initialize conversation chain
#     conversational_chain = get_conversation_chain(retriever)
#     st.session_state.conversational_chain = conversational_chain

#     # Initialize chat history if not already in session state
#     if 'chat_history' not in st.session_state:
#         st.session_state.chat_history = []
# except Exception as e:
#     st.error(f"Error preparing documents: {e}")

# # Chat Input Section
# user_input = st.text_input("Ask a question about the documents:")

# # Buttons for submitting and managing chat
# col1, col2 = st.columns([1, 1])
# with col1:
#     if st.button("Submit"):
#         if user_input:
#             try:
#                 # Pass the user input through the conversation chain
#                 response = st.session_state.conversational_chain(
#                     {"question": user_input, "chat_history": st.session_state.chat_history}
#                 )

#                 # Display response
#                 st.session_state.chat_history.append({"question": user_input, "answer": response['answer']})
#                 st.write(f"**Bot:** {response['answer']}")
#             except Exception as e:
#                 st.error(f"Error generating response: {e}")
#         else:
#             st.warning("Please enter a question!")

# with col2:
#     if st.button("Clear Chat"):
#         st.session_state.chat_history = []
#         st.success("Chat history cleared!")

# # if st.session_state.chat_history:
# #     for message in st.session_state.chat_history:
# #         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
# #         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

#         if message.get("context_docs"):
#             with st.expander("Source Documents"):
#                 for doc in message["context_docs"]:
#                     st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
#                     st.write(doc.page_content)

#================================code with Next quetion button==============
# import streamlit as st
# from langchain.document_loaders import DirectoryLoader, PyPDFLoader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_community.llms import Ollama
# from langchain.vectorstores import FAISS
# from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
# from langchain_core.runnables.history import RunnableWithMessageHistory
# from langchain_community.chat_message_histories import ChatMessageHistory
# from langchain.chains import create_retrieval_chain
# from langchain.chains.combine_documents import create_stuff_documents_chain
# from sentence_transformers import SentenceTransformer, util
# from langchain_core.chat_history import BaseChatMessageHistory
# from langchain.chains import create_history_aware_retriever
# from langchain.embeddings import HuggingFaceEmbeddings
# import os

# # Define template for bot and user messages
# bot_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px;">
#     <div style="flex-shrink: 0; margin-right: 10px;">
#         <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>
#     <div style="background-color: #f1f1f1; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

# user_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px; justify-content: flex-end;">
#     <div style="flex-shrink: 0; margin-left: 10px;">
#         <img src="https://cdn.iconscout.com/icon/free/png-512/free-q-characters-character-alphabet-letter-36051.png?f=webp&w=512" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>    
#     <div style="background-color: #007bff; color: white; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

# # Function to prepare and split documents from a specified directory
# def prepare_and_split_docs(pdf_directory):
#     loader = DirectoryLoader(pdf_directory, glob="**/*.pdf", loader_cls=PyPDFLoader)
#     documents = loader.load()
#     splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
#         chunk_size=512,
#         chunk_overlap=256,
#         disallowed_special=(),
#         separators=["\n\n", "\n", " "]
#     )
#     return splitter.split_documents(documents)

# # Function to ingest documents into the vector database
# def ingest_into_vectordb(split_docs):
#     embeddings = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
#     return FAISS.from_documents(split_docs, embeddings)

# def get_conversation_chain(retriever):
#     llm = Ollama(model="meta-llama/Llama-3.2-3B")
#     contextualize_q_system_prompt = (
#         "Given the chat history and the latest user question, "
#         "provide a response that directly addresses the user's query based on the provided documents. "
#         "If no relevant answer is found, respond with: "
#         "'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?' "
#         "Do not rephrase the question or ask follow-up questions."
#     )
#     contextualize_q_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", contextualize_q_system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )
#     history_aware_retriever = create_history_aware_retriever(
#         llm, retriever, contextualize_q_prompt
#     )
#     system_prompt = (
#         "As a chat assistant, provide accurate and relevant information based on the provided document in 2-3 sentences. "
#         "Answer should be limited to 50 words and 2-3 sentences. If no relevant information is found, respond with: "
#         "'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?' "
#         "{context}"
#     )
#     qa_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )
#     question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)
#     rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)

#     store = {}

#     def get_session_history(session_id: str) -> BaseChatMessageHistory:
#         if session_id not in store:
#             store[session_id] = ChatMessageHistory()
#         return store[session_id]

#     return RunnableWithMessageHistory(
#         rag_chain,
#         get_session_history,
#         input_messages_key="input",
#         history_messages_key="chat_history",
#         output_messages_key="answer",
#     )

# # Main Streamlit app
# logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png"
# st.image(logopath, width=200)
# st.title("AssureBot")
# st.write("Welcome! How can I help you today?")

# # Define the path to your PDF directory
# pdf_directory = "Data/SalesDocs"  # Update with your actual PDF folder path

# # Prepare and ingest documents into the vector store
# split_docs = prepare_and_split_docs(pdf_directory)
# vector_db = ingest_into_vectordb(split_docs)
# retriever = vector_db.as_retriever()

# # Initialize the conversation chain
# if 'conversational_chain' not in st.session_state:
#     st.session_state.conversational_chain = get_conversation_chain(retriever)

# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []

# if 'user_input' not in st.session_state:
#     st.session_state.user_input = ""

# # Chat input
# user_input = st.text_input("Ask a question about the documents:", value=st.session_state.user_input)

# # Buttons for submitting and refreshing
# # Buttons for submitting and refreshing
# col1, col2 = st.columns([1, 1])
# with col1:
#     if st.button("Submit"):
#         user_input_lower = user_input.lower()  # Normalize input for checking

#         # Store the user input in session state
#         st.session_state.user_input = user_input

#         # Respond to greetings
#         if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
#             response = "Hello! How can I assist you?"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to thanks
#         elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
#             response = "Thank you! Let me know if you have any more queries."
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to yes/no
#         elif any(yn in user_input_lower for yn in ["yes", "no"]):
#             response = "I understand. If you have any other questions, feel free to ask!"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Process document queries
#         elif user_input and 'conversational_chain' in st.session_state:
#             session_id = "abc123"  # Static session ID for this demo; you can make it dynamic if needed
#             response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#             context_docs = response.get('context', [])

#             # Check if the response is the default message
#             if response['answer'] == "I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?":
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": []})
#             else:
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": context_docs})

# with col2:
#     if st.button("Next Question"):
#         st.session_state.user_input = ""  # Clear the input without refreshing the app

# # Display chat history
# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

#         # Display the source documents if available
#         if message.get('context_docs'):
#             with st.expander("Source Documents"):
#                 for doc in message['context_docs']:
#                     st.write(f"Source: {doc.metadata['source']}")
#                     st.write(doc.page_content)
# import streamlit as st
# from langchain.document_loaders import DirectoryLoader, PyPDFLoader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_community.llms import Ollama
# from langchain.vectorstores import FAISS
# from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
# from langchain_core.runnables.history import RunnableWithMessageHistory
# from langchain_community.chat_message_histories import ChatMessageHistory
# from langchain.chains import create_retrieval_chain
# from langchain.chains.combine_documents import create_stuff_documents_chain
# from sentence_transformers import SentenceTransformer, util
# from langchain_core.chat_history import BaseChatMessageHistory
# from langchain.chains import create_history_aware_retriever
# from langchain.embeddings import HuggingFaceEmbeddings
# import os

# # Define template for bot and user messages
# bot_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px;">
#     <div style="flex-shrink: 0; margin-right: 10px;">
#         <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>
#     <div style="background-color: #f1f1f1; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

# user_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px; justify-content: flex-end;">
#     <div style="flex-shrink: 0; margin-left: 10px;">
#         <img src="https://cdn.iconscout.com/icon/free/png-512/free-q-characters-character-alphabet-letter-36051.png?f=webp&w=512" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>    
#     <div style="background-color: #007bff; color: white; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

# # Function to prepare and split documents from a specified directory
# def prepare_and_split_docs(pdf_directory):
#     loader = DirectoryLoader(pdf_directory, glob="**/*.pdf", loader_cls=PyPDFLoader)
#     documents = loader.load()
#     splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
#         chunk_size=512,
#         chunk_overlap=256,
#         disallowed_special=(),
#         separators=["\n\n", "\n", " "]
#     )
#     return splitter.split_documents(documents)

# # Function to ingest documents into the vector database
# def ingest_into_vectordb(split_docs):
#     embeddings = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
#     return FAISS.from_documents(split_docs, embeddings)

# def get_conversation_chain(retriever):
#     llm = Ollama(model="llama3.2")
#     contextualize_q_system_prompt = (
#         "Given the chat history and the latest user question, "
#         "provide a response that directly addresses the user's query based on the provided documents. "
#         "If no relevant answer is found, respond with: "
#         "'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?' "
#         "Do not rephrase the question or ask follow-up questions."
#     )
#     contextualize_q_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", contextualize_q_system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )
#     history_aware_retriever = create_history_aware_retriever(
#         llm, retriever, contextualize_q_prompt
#     )
#     system_prompt = (
#         "As a chat assistant, provide accurate and relevant information based on the provided document in 2-3 sentences. "
#         "Answer should be contextualize and relevant words and  sentences. If no relevant information is found, respond with: "
#         "'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?' "
#         "{context}"
#     )
#     qa_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )
#     question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)
#     rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)

#     store = {}

#     def get_session_history(session_id: str) -> BaseChatMessageHistory:
#         if session_id not in store:
#             store[session_id] = ChatMessageHistory()
#         return store[session_id]

#     return RunnableWithMessageHistory(
#         rag_chain,
#         get_session_history,
#         input_messages_key="input",
#         history_messages_key="chat_history",
#         output_messages_key="answer",
#     )

# # Main Streamlit app
# st.set_page_config(page_title="AssureBot", layout="wide")
# logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png"
# st.image(logopath, width=200)
# st.title("AssureBot")
# st.write("Welcome! How can I help you today?")

# # Define the path to your PDF directory
# pdf_directory = "Data/SalesDocs"  # Update with your actual PDF folder path

# # Prepare and ingest documents into the vector store
# split_docs = prepare_and_split_docs(pdf_directory)
# vector_db = ingest_into_vectordb(split_docs)
# retriever = vector_db.as_retriever()

# # Initialize the conversation chain
# if 'conversational_chain' not in st.session_state:
#     st.session_state.conversational_chain = get_conversation_chain(retriever)

# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []

# if 'user_input' not in st.session_state:
#     st.session_state.user_input = ""

# # Chat input
# user_input = st.text_input("Ask a question about the documents:", value=st.session_state.user_input)

# # Buttons for submitting and next question
# col1, col2 = st.columns([1, 1])
# with col1:
#     if st.button("Submit"):
#         user_input_lower = user_input.lower()  # Normalize input for checking

#         # Store the user input in session state
#         st.session_state.user_input = user_input

#         # Respond to greetings
#         if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
#             response = "Hello! How can I assist you?"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to thanks
#         elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
#             response = "Thank you! Let me know if you have any more queries."
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to yes/no
#         elif any(yn in user_input_lower for yn in ["yes", "no"]):
#             response = "I understand. If you have any other questions, feel free to ask!"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Process document queries
#         elif user_input and 'conversational_chain' in st.session_state:
#             session_id = "abc123"  # Static session ID for this demo; you can make it dynamic if needed
#             response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#             context_docs = response.get('context', [])

#             # Check if the response is the default message
#             if response['answer'] == "I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?":
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": []})
#             else:
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": context_docs})

# with col2:
#     if st.button("Next Question"):
#         st.session_state.user_input = ""  # Clear the input

# # Display chat history
# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

#         # Display the source documents if available
#         if message.get('context_docs'):
#             with st.expander("Source Documents"):
#                 for doc in message['context_docs']:
#                     st.write(f"Source: {doc.metadata['source']}")
#                     st.write(doc.page_content)
# import streamlit as st
# from langchain.document_loaders import DirectoryLoader, PyPDFLoader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_community.llms import Ollama
# from langchain.vectorstores import FAISS
# from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
# from langchain_core.runnables.history import RunnableWithMessageHistory
# from langchain_community.chat_message_histories import ChatMessageHistory
# from langchain.chains import create_retrieval_chain
# from langchain.chains.combine_documents import create_stuff_documents_chain
# from sentence_transformers import SentenceTransformer, util
# from langchain_core.chat_history import BaseChatMessageHistory
# from langchain.chains import create_history_aware_retriever
# from langchain.embeddings import HuggingFaceEmbeddings


# # Templates for user and bot messages
# bot_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px;">
#     <div style="flex-shrink: 0; margin-right: 10px;">
#         <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>
#     <div style="background-color: #f1f1f1; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

# user_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px; justify-content: flex-end;">
#     <div style="flex-shrink: 0; margin-left: 10px;">
#         <img src="https://cdn.iconscout.com/icon/free/png-512/free-q-characters-character-alphabet-letter-36051.png?f=webp&w=512" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>    
#     <div style="background-color: #007bff; color: white; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

# # Function to prepare and split documents
# def prepare_and_split_docs(pdf_directory):
#     loader = DirectoryLoader(pdf_directory, glob="**/*.pdf", loader_cls=PyPDFLoader)
#     documents = loader.load()
#     splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
#         chunk_size=512,
#         chunk_overlap=256,
#         disallowed_special=(),
#         separators=["\n\n", "\n", " "]
#     )
#     split_docs = splitter.split_documents(documents)
#     return split_docs

# # Function to ingest documents into vector store
# def ingest_into_vectordb(split_docs):
#     embeddings = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
#     db = FAISS.from_documents(split_docs, embeddings)
#     return db

# # Function to create a conversation chain
# def get_conversation_chain(retriever):
#     llm = Ollama(model="llama3.2")
#     contextualize_q_system_prompt = (
#         "Given the chat history and the latest user question, "
#         "provide a response based on the documents. If no answer is found, "
#         "respond: 'I'm sorry, but I couldn’t find an answer. Could you rephrase or provide more details?'"
#     )
#     contextualize_q_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", contextualize_q_system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )
#     history_aware_retriever = create_history_aware_retriever(
#         llm, retriever, contextualize_q_prompt
#     )
#     system_prompt = (
#         "As a chat assistant, provide accurate and relevant information based on the provided document in 2-3 sentences. "
#         "Answer should be correct to the point short and brief for given question. If no relevant information is found, respond with: "
#         "'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?' "
#         "{context}"
#     )
#     qa_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )
#     question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)
#     rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)

#     store = {}

#     def get_session_history(session_id: str):
#         if session_id not in store:
#             store[session_id] = ChatMessageHistory()
#         return store[session_id]

#     conversational_rag_chain = RunnableWithMessageHistory(
#         rag_chain,
#         get_session_history,
#         input_messages_key="input",
#         history_messages_key="chat_history",
#         output_messages_key="answer",
#     )
#     return conversational_rag_chain

# # Main Streamlit app
# logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png"
# st.image(logopath, width=200)
# st.title("Q&A Bot")
# st.write("Welcome! How can I help you today?")

# pdf_directory = "Data/SalesDocs"  # Update to your actual folder path

# # Prepare documents and ingest into vector store
# split_docs = prepare_and_split_docs(pdf_directory)
# vector_db = ingest_into_vectordb(split_docs)
# retriever = vector_db.as_retriever()

# # Initialize conversation chain
# conversational_chain = get_conversation_chain(retriever)
# st.session_state.conversational_chain = conversational_chain

# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []

# # Chat input
# user_input = st.text_input("Ask a question about the documents:")
# # Buttons for submitting and refreshing
# col1, col2 = st.columns([1, 1])
# with col1:
#     if st.button("Submit"):
#         user_input_lower = user_input.lower()  # Normalize input for checking

#         # Store the user input in session state
#         st.session_state.user_input = user_input

#         # Respond to greetings
#         if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
#             response = "Hello! How can I assist you?"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to thanks
#         elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
#             response = "Thank you! Let me know if you have any more queries."
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to yes/no
#         elif any(yn in user_input_lower for yn in ["yes", "no"]):
#             response = "I understand. If you have any other questions, feel free to ask!"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Process document queries
#         elif user_input and 'conversational_chain' in st.session_state:
#             session_id = "abc123"  # Static session ID for this demo; you can make it dynamic if needed
#             response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#             context_docs = response.get('context', [])

#             # Check if the response is the default message
#             if response['answer'] == "I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?":
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": []})
#             else:
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": context_docs})

# with col2:
#     if st.button("Refresh"):
#         st.session_state.user_input = ""  # Clear the input
#         st.session_state.chat_history = []  # Clear chat history

# # Display chat history
# for message in st.session_state.chat_history:
#     st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#     st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

#     if message.get("context_docs"):
#         with st.expander("Source Documents"):
#             for doc in message["context_docs"]:
#                 st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
#                 st.write(doc.page_content)
